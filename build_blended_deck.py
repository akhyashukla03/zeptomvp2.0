import os
import sys

# Force stdout UTF-8 encoding for Windows terminal printing
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def build_pptx_deck():
    prs = Presentation()
    # Widescreen 16:9 (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Ultra-Clean Executive Theme Palette (Light Canvas, High Contrast, Colorblind Safe)
    BG_LIGHT = RGBColor(248, 250, 252)       # #f8fafc Light Slate Canvas
    CARD_BG = RGBColor(255, 255, 255)        # #ffffff Crisp Pure White Card
    CARD_BORDER = RGBColor(203, 213, 225)    # #cbd5e1 Subtle Border
    HEADER_GRAD = RGBColor(59, 7, 100)        # #3b0764 Zepto Deep Violet Header
    PRIMARY_TITLE = RGBColor(15, 23, 42)      # #0f172a Deep Slate Black
    SUBTITLE_COLOR = RGBColor(79, 70, 229)    # #4f46e5 Indigo Accent
    BOX_HEADER_COLOR = RGBColor(67, 56, 202)  # #4338ca Royal Indigo Header
    BODY_TEXT_COLOR = RGBColor(30, 41, 59)    # #1e293b Dark Slate Body Text
    WHITE = RGBColor(255, 255, 255)
    MUTED_TEXT = RGBColor(100, 116, 139)     # #64748b Slate Muted Text
    CALLOUT_BG = RGBColor(241, 245, 249)     # #f1f5f9 Soft Callout Tint
    CALLOUT_BORDER = RGBColor(79, 70, 229)   # #4f46e5 Indigo Border
    GREEN_ACCENT = RGBColor(4, 120, 87)      # #047857 Emerald Green Accent

    steps_list = ["Context", "Market", "Research", "Insights", "Canvas", "Ideation", "MVP", "Architecture", "Metrics", "GTM"]

    # Image paths (100% Crisp Vector SVG Renders)
    img_dir = os.path.join(os.path.dirname(__file__), 'images', 'crisp_vector_renders')
    img_fallback = os.path.join(os.path.dirname(__file__), 'images', 'tight_iphone_app.png')
    
    def get_img_path(fname):
        p = os.path.join(img_dir, fname)
        return p if os.path.exists(p) else img_fallback

    img1 = get_img_path('01_b2b_free_sampler_cart.png')
    img2 = get_img_path('02_category_streak_quest_board.png')
    img3 = get_img_path('03_restock_cadence_interception.png')
    img4 = get_img_path('04_post_trial_outcome_loop.png')
    img5 = get_img_path('05_lifecycle_moment_interceptor.png')
    img6 = get_img_path('06_freshness_guaranteed_audit.png')
    img7 = get_img_path('07_neighbourhood_trend_feed.png')
    img8 = get_img_path('08_discovery100_voucher_conversion.png')

    figma_url = "https://www.figma.com/design/2gZJHtjmpnI677IhC66PqZ/Untitled?node-id=0-1"
    app_url = "https://zeptomvp2.streamlit.app"

    slides_data = [
        {
            "slide_num": 1,
            "tagline": "CATEGORY STREAK + GROCERY POINTS CROSS-SUBSIDY MODEL",
            "title": "1. Zepto Expands Blended Gross Margin (+300bps) by Converting Daily Grocery Habits into Recurring Multi-Category Buying",
            "subtitle": "Funding grocery point rewards from high-margin non-grocery categories to bypass traditional P&L subsidy limits.",
            "box1_title": "📋 STRATEGIC BRIEF & PM SCOPE",
            "box1_bullets": [
                "Role & Scope: PM on Zepto Growth Team driving trust-led cross-category discovery.",
                "Routine Tunnel Vision: 71.2% of active users are locked into fast daily grocery staple checkouts (<45 seconds).",
                "Strategic Objective: Lift Monthly Active Customers (MAC) buying from 2+ categories from 8.2% to 28.4% in 12 months.",
                "Target Categories: Personal Care (35% margin), Pet Supplies (45% margin), Electronics (30% margin), Baby Care (35% margin)."
            ],
            "box2_title": "⚡ THE CROSS-SUBSIDY ARBITRAGE MODEL",
            "box2_bullets": [
                "Grocery Point Cross-Subsidy: Earning 2x points on morning staples (milk, eggs) offsets the customer's grocery bill.",
                "Sustainable Arbitrage: Points are funded from high non-grocery margins (~35–45%), not from direct P&L discounts.",
                "Ecosystem Lock-in: Competitors cannot copy this cross-subsidy without gross margin arbitrage (or subsidizing from P&L).",
                "Compounding Flywheel: More category exploration -> more grocery points -> stronger grocery habit retention."
            ],
            "image_path": img1,
            "bottom_title": "🔗 PORTAL DIRECTORY & VERIFIED HYPERLINKS",
            "bottom_text": f"Live Streamlit App: {app_url} | Figma Design System: {figma_url}",
            "app_url": app_url,
            "figma_url": figma_url
        },
        {
            "slide_num": 2,
            "tagline": "CONVENIENCE-JUSTIFIED TOP-UPS VS PLANNED BULK BUYS",
            "title": "2. Users Leak Planned Non-Grocery Spending to Amazon/DMart; MVP Targets Restock Cadence and Trial-to-Repeat Loops",
            "subtitle": "Shifting focus from massive bulk orders to convenience-justified trials and replenishment intercepts.",
            "box1_title": "⭐ WHAT THE ECOSYSTEM HAS SOLVED",
            "box1_bullets": [
                "10-Minute Hyperlocal Delivery Speed (Zepto, Blinkit, Instamart).",
                "Dark Store Density & Real-Time Stock (500+ Dark Stores).",
                "Sub-Second Cart Search & Auto-Complete UI."
            ],
            "box2_title": "❌ CADENCE & TRIAL-TO-REPEAT BLIND SPOTS (N=24 DATA)",
            "box2_bullets": [
                "Purchasing Leakage: 83% (20/24) buy planned non-grocery items on Amazon/Nykaa/DMart [Survey Q3, N=24].",
                "Planned Mismatch: 33% (8/24) cite planned bulk-buy mismatch as the #1 barrier [Survey Q4].",
                "Zero Cadence Nudges: No Q-Commerce platform currently offers reorder suggestions based on user time gaps.",
                "Trust & Quality Friction: 20.1% fear quality degradation; 17% (4/24) fear refund chatbot loops."
            ],
            "image_path": img6,
            "bottom_title": "★ CORE THESIS TESTED & VALIDATED",
            "bottom_text": "Exploration stalls because trial samples are treated as one-offs. Fix: Free sample -> Category Streak Quest -> grocery offset. Restock cadence nudge intercepts Amazon orders 2 days before purchase.",
            "app_url": app_url,
            "figma_url": figma_url
        },
        {
            "slide_num": 3,
            "tagline": "AI REVIEW WORKFLOW & PRIMARY RESEARCH SYNTHESIS",
            "title": "3. Multi-Platform Mixed Review Corpus (10k Items) & Survey Isolate 3 Core Pillars: Delights, Feature Requests & Barriers",
            "subtitle": "Synthesizing 10,000 mixed reviews (40% Positive, 30% Neutral, 30% Negative) and primary N=24 survey into actionable insights.",
            "box1_title": "📊 AI PIPELINE & SYNTHETIC CORPUS (10,000 ITEMS)",
            "box1_bullets": [
                "Pipeline: Ingested 10,000 records across 10 channels -> 40% Positive Delights, 30% Neutral Requests, 30% Frictions.",
                "40% Positive Delights (4,000 items): High adoption for sub-10 min emergency refills (sunscreen, travel grooming, pet treats).",
                "30% Neutral Feature Requests (3,000 items): Unmet demand for Shade Matchers, Doorstep Try & Inspect, and Category Streaks."
            ],
            "box2_title": "🎯 3 SYNTHESIZED BARRIER BUCKETS (N=24 SURVEY)",
            "box2_bullets": [
                "1. Mismatch Barrier (33%, 8/24): Planned bulk buys elsewhere — intercept 2 days prior via cadence data.",
                "2. Trial-to-Repeat Gap: Sampling exists but lacks an outcome loop to convert initial trial to repeat buying.",
                "3. Cross-Subsidy Appeal (63%, 15/24): Prefer grocery points on streaks over flat cashbacks (38% trial via sample; 92% trust doorstep swap)."
            ],
            "image_path": img8,
            "bottom_title": "💡 RESEARCH CONCLUSION & INSIGHT SYNTHESIS",
            "bottom_text": "Users do not need generic ads. They need Category Streaks (lock-in loyalty), Cadence Nudges (intercept bulk buys), and Outcome Loop Cards (drive repeat).",
            "app_url": app_url,
            "figma_url": figma_url
        },
        {
            "slide_num": 4,
            "tagline": "HIGH-FREQUENCY GROCERY BUYERS WHO LEAK NON-GROCERY SPEND",
            "title": "4. High-Frequency Grocery Refillers (79% Weekly+) Form the Prime Wedge to Intercept Planned Non-Grocery Restocks",
            "subtitle": "Targeting retained grocery refuelers and using grocery cadence to intercept Amazon purchases 2 days before restock.",
            "box1_title": "📊 TARGET COHORT PROFILE & RATIONALE",
            "box1_bullets": [
                "Cohort Definition: Weekly+ grocery users (79%, 19/24) leaking non-grocery spend (83%, 20/24) [Survey N=24].",
                "Strategic Fit: Retained, high-frequency users yield lower CAC than acquiring new users.",
                "Predictive Interception: Grocery cadence predicts non-grocery restock timing -> nudge fires 2 days before Amazon order.",
                "Trust Transfer: Daily Zepto staple use lowers friction for adjacent category trials."
            ],
            "box2_title": "👤 PERSONA: NEHA — SKINCARE FAN, 26, BANGALORE",
            "box2_bullets": [
                "JTBD: When ordering morning milk, I want to try skincare samples so discovery offsets my grocery bill.",
                "Quote: \"If trying a new brand gets me 2x points on daily bread, it becomes a monthly habit.\"",
                "Pain: Cautious Explorer — wants to try beauty care but fears dark-store expiry and refund friction.",
                "Fit: Risk-free ₹0 sample + streak point converts existing demand with zero operational overhaul."
            ],
            "image_path": img3,
            "bottom_title": "🎯 IMPACT IF SOLVED FOR THE CAUTIOUS EXPLORER",
            "bottom_text": "Cuts pre-checkout anxiety loops. Makes multi-category exploration safe. Boosts customer LTV up to 3.4x via B2B sampling and loyalty points lock-in.",
            "app_url": app_url,
            "figma_url": figma_url
        },
        {
            "slide_num": 5,
            "tagline": "SUPPLY-SIDE MONETIZATION & UNIT ECONOMICS RIGOR",
            "title": "5. Brand-Funded Cross-Subsidy Yields Positive Unit Economics with Immediate Payback on Cross-Category LTV",
            "subtitle": "Grocery rewards are funded by non-grocery margins (~35–45%), not from Zepto P&L.",
            "box1_title": "📈 TAM / SAM / SOM OPPORTUNITY SIZING",
            "box1_bullets": [
                "TAM: $18.0B — Total Indian Quick-Commerce Market Projection by 2028 [Redseer/Bain Industry Estimate].",
                "SAM: $4.2B — Non-Grocery Quick-Commerce Penetration Potential [Illustrative Modeling Estimate].",
                "SOM: $480M — Zepto Cross-Category Discovery Capture Opportunity [Illustrative Modeling Estimate]."
            ],
            "box2_title": "💰 UNIT ECONOMICS & CROSS-SUBSIDY MATH",
            "box2_bullets": [
                "B2B Brand Sample Fee: Brand pays ₹15/unit -> Zepto earns ₹15 at zero extra delivery cost.",
                "Streak Reward: 2x points on ~₹50 grocery order ≈ ₹5 reward; funded from non-grocery margin.",
                "Cross-Subsidy Math: ₹499 BPC order @ 35% margin = ₹175 gross profit; easily covers ₹5 streak reward x 35 orders.",
                "Net Impact: Cross-category trial is margin-accretive at scale (+300bps gross margin lift) — not a discount or subsidy."
            ],
            "image_path": None,
            "bottom_title": "💡 FINANCIAL FLYWHEEL SUMMARY",
            "bottom_text": "Shifting grocery refillers to 35%–50% margin categories expands blended gross margin by +300bps and captures $480M SOM.",
            "app_url": app_url,
            "figma_url": figma_url
        },
        {
            "slide_num": 6,
            "tagline": "COMPOUNDING DATA FLYWHEEL VS COPIABLE PROMOS",
            "title": "6. Zepto's Defensibility Lies in its Restock Cadence Engine and Local Purchase-Graph, Not Copiable Promos",
            "subtitle": "Flat discounts are easily copied; restock history and PIN-code purchase graphs create years of defensibility.",
            "box1_title": "⚡ TRIVIALLY COPIABLE VS COMPOUNDING MOATS",
            "box1_bullets": [
                "Copiable Promos: Free samples/coupons and 15-min swap policies can be copied by Blinkit/Instamart in a sprint.",
                "Moat 1 — Cadence Engine: 6–12 months of per-user order history required. Rivals starting today face a 1-year data gap.",
                "Moat 2 — PIN-Code Social Graph: Years of dark-store density. Amazon/DMart structurally cannot replicate.",
                "Moat 3 — Outcome Attribution: Trial-to-repeat data -> Zepto Atom. No banner ad can match this performance."
            ],
            "box2_title": "🏆 THREE STRATEGIC HORIZONS (RICE MATRIX)",
            "box2_bullets": [
                "Horizon 1 (MVP) — Streak + Cadence + Outcome: Zero CapEx (brands fund samples). RICE: 230 [Vetted MVP].",
                "Horizon 2 (Growth) — Lifecycle + Neighbourhood Feed: Self-declared life moments + opt-in trend feed (50+ user min). RICE: 190.",
                "Horizon 3 (Vision) — Predictive Auto-Replenishment: Set-and-forget cross-category companion. RICE: 150."
            ],
            "image_path": None,
            "bottom_title": "💡 WHY HORIZON 1 WINS FIRST",
            "bottom_text": "Horizon 1 wins first: zero CapEx, attacks the #1 barrier directly, and embeds in recurring grocery bags. H2/H3 require the data H1 generates.",
            "app_url": app_url,
            "figma_url": figma_url
        },
        {
            "slide_num": 7,
            "tagline": "THE FLYWHEEL MVP: STREAKS, CADENCE NUDGES, OUTCOME CARDS",
            "title": "7. The Discovery MVP Embeds Category Streaks, Cadence Nudges, and Post-Trial Outcome Loop Cards",
            "subtitle": "Connecting trials, streaks, and intercepts into a closed-loop category adoption system.",
            "box1_title": "💎 3 CORE BUILT MVP CAPABILITIES",
            "box1_bullets": [
                "1. 🏆 Category Streak Board: Pre-stamped Pantry sticker (Endowed Progress Effect: 34% vs 19% completion rate) -> 2x grocery points.",
                "2. 🕒 Cadence Interception Nudge: Restock timing engine fires 2 days before Amazon bulk buy (offers ₹0 brand-funded trial pack).",
                "3. 📦 Outcome Loop Card: Trial-to-repeat converter fires at 48h with PIN-code proof ('847 near you reordered'). Rate & Add to next cart."
            ],
            "box2_title": "🏆 HORIZON 2 EXPANSIONS & TRUST SIGNALS",
            "box2_bullets": [
                "4. 🔄 Lifecycle Moment Interceptor: Basket signal prompts self-declaration card (baby diaper addition).",
                "5. 🔥 Neighbourhood Trend Feed: Opt-in, 50+ user min, sensitive-category blocklist, PIN-code dark-store data moat.",
                "6. 🛡️ Freshness Guaranteed Badge: Checkout trust signal — zero IoT hardware infrastructure required at MVP stage."
            ],
            "image_path": img2,
            "bottom_title": "🎨 LIVE FIGMA DESIGN SYSTEM & INTERACTIVE MVP PORTAL",
            "bottom_text": f"Live MVP Portal: {app_url} | Figma Design System: {figma_url}",
            "app_url": app_url,
            "figma_url": figma_url
        },
        {
            "slide_num": 8,
            "tagline": "FROM A TYPED GROCERY CART TO A TRUSTED CROSS-CATEGORY ORDER. SAME APP, TWO VIEWS.",
            "title": "8. Four-Layer Decision Engine Powers Predictive Restock Nudges and Closed-Loop Brand Attribution",
            "subtitle": "Four-layer architecture integrates brand inventory, cart triggers, and loyalty ledgers.",
            "box1_title": "⚙ SYSTEM ARCHITECTURE (4 CORE LAYERS)",
            "box1_bullets": [
                "Layer 1 (Client UI): Mobile Cart UI + pre-stamped Streak Board + Cadence widget + Outcome Card + B2B Sampler Carousel.",
                "Layer 2 (Decision Engine): Restock Cadence Engine (6–12mo history) + Endowed Progress Tracker + Social Graph Parser.",
                "Layer 3 (Operations): Dark-store picker checklist (sample insertion <3s SLA) + Freshness Guaranteed badge.",
                "Layer 4 (B2B Marketplace): Brand attribution dashboard (sample-to-repeat conversion rates per cohort via Zepto Atom)."
            ],
            "box2_title": "🔄 USER EMOTION & METRIC MAPPING ACROSS 4 STAGES",
            "box2_bullets": [
                "Stage 1 (Cart - Curious): Types grocery staples -> streak board + cadence signal flagged.",
                "Stage 2 (Nudge - Confident): Claims ₹0 sample -> packed in bag in <3s with zero checkout friction.",
                "Stage 3 (Follow-Up - Reassured): Outcome Loop Card at 48h -> sees '847 near you reordered' + streak CTA.",
                "Stage 4 (Repeat - Empowered): Rates + adds full-size -> streak advances -> 2x grocery points earned."
            ],
            "image_path": img4,
            "bottom_title": "💡 TECHNICAL ARCHITECTURE MOAT & FEASIBILITY",
            "bottom_text": "Zero hardware cost for MVP sampler. H2 optional S3 photo storage capped at <$15/month via 7-day TTL lifecycle rules.",
            "app_url": app_url,
            "figma_url": figma_url
        },
        {
            "slide_num": 9,
            "tagline": "INCREMENTAL LIFT VIA RANDOMIZED HOLDOUT GROUPS",
            "title": "9. MCER Growth (8.2% -> 28.4%) Is Validated via Randomized Holdout Groups and SLA Guardrails",
            "subtitle": "Measuring true incremental lift via control groups rather than raw attach proxies.",
            "box1_title": "⭐ NORTH STAR & INCREMENTALITY DESIGN",
            "box1_bullets": [
                "North Star Metric: Monthly Category Exploration Rate (MCER) — % of MACs purchasing from 2+ categories/month.",
                "Trajectory Target: Baseline: 8.2% -> M3 (Pilot): 14.5% -> M6 (Metro): 21.0% -> M12 Target: 28.4% [Illustrative].",
                "Holdout Control: 10% randomized control group (no sample/streak/nudge); MCER lift = treatment delta vs control.",
                "Leading Indicators: Sample attach rate (>35%), trial-to-repeat rate (≥12%, 14 days), streak completion rate (≥40%)."
            ],
            "box2_title": "🛡️ OPERATIONAL GUARDRAILS & KILL THRESHOLDS",
            "box2_bullets": [
                "Picker SLA Floor: Packing time addition capped at <3 seconds per order. Checkout drop-off cap <0.2%.",
                "Refund Rate Ceiling: Doorstep replacements below 1.5% of non-grocery orders. Privacy Floor: 50+ users/PIN min.",
                "Kill Threshold 1: Sample attach < 15% after 30 days -> stop sampling, pivot to nudge-only.",
                "Kill Threshold 2: Trial-to-repeat < 8% in 14 days -> retrain model. MCER lift < 2pp at 90 days -> pause and redesign."
            ],
            "image_path": img7,
            "bottom_title": "📈 METRIC COMPOUNDING & INTEGRITY",
            "bottom_text": "MCER is measured strictly from real event streams: sample claims, streak completions, voucher redemptions. Zero proxies. Holdout ensures causality.",
            "app_url": app_url,
            "figma_url": figma_url
        },
        {
            "slide_num": 10,
            "tagline": "HORIZON 2 ROADMAP & PRIVACY GUARDRAILS",
            "title": "10. Phased GTM Roadmap Launches Horizon 1 MVP While Building Horizon 2 Privacy-Guardrailed Features",
            "subtitle": "Phased rollout for Category Streaks, Cadence Intercepts, Lifecycle Moments, and Trend Feeds.",
            "box1_title": "🚀 4-PHASE GTM ROLLOUT STRATEGY",
            "box1_bullets": [
                "Phase 1 (30 Days - MVP Beta Bangalore): Streak + Cadence Nudge + Outcome Cards across 10 dark stores with Cetaphil & Pedigree.",
                "Phase 2 (90 Days - Metro Rollout): Expand H1 to all Tier-1 metros (150 dark stores); launch B2B brand portal (Gate: MCER lift ≥3pp).",
                "Phase 3 (180 Days - Horizon 2 Beta): Lifecycle Interceptor (self-declared) + Neighbourhood Feed (50+ user min, opt-in).",
                "Phase 4 (GA - National Scale): Full rollout across 500+ dark-store hubs. H3 auto-replenishment roadmap."
            ],
            "box2_title": "⚠️ RISKS, MITIGATIONS & DIRECTORY",
            "box2_bullets": [
                "Rival Copying -> Streak + grocery cross-subsidy requires rivals to have identical non-grocery margin structure.",
                "Data Runway -> Start cadence data collection Day 1 of pilot; 6–12 months needed for accurate interception.",
                "Privacy Backlash -> Opt-in only, 50+ user PIN-code min, sensitive-category blocklist, DPDP-compliant consent.",
                f"Live MVP Portal: {app_url} | Figma Design System: {figma_url}"
            ],
            "image_path": img5,
            "bottom_title": "🔒 PRIVACY PRINCIPLE & DIRECTORY",
            "bottom_text": f"Privacy principle: Zepto knows what you buy. It should never make you feel watched. Live Demo: {app_url} | Figma: {figma_url}",
            "app_url": app_url,
            "figma_url": figma_url
        }
    ]

    for data in slides_data:
        slide = prs.slides.add_slide(blank_layout)

        # 1. Background Fill (Light Slate Canvas)
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = BG_LIGHT
        bg_shape.line.color.rgb = BG_LIGHT

        # 2. Top Header Tagline Banner (Deep Violet Pill Banner)
        banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.35), Inches(12.333), Inches(0.45))
        banner.fill.solid()
        banner.fill.fore_color.rgb = HEADER_GRAD
        banner.line.color.rgb = HEADER_GRAD
        tf_b = banner.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = Inches(0.18)
        tf_b.margin_top = Inches(0.08)
        p_b = tf_b.paragraphs[0]
        p_b.text = f"  {data['tagline'].upper()}"
        p_b.font.size = Pt(14)  # STRICT MINIMUM 14PT
        p_b.font.bold = True
        p_b.font.color.rgb = WHITE

        # 3. Slide Title & Subtitle (Deep Slate Black & Indigo Accent)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.88), Inches(12.333), Inches(1.05))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = Inches(0)
        tf_t.margin_top = Inches(0)
        
        p_t = tf_t.paragraphs[0]
        p_t.text = data["title"]
        p_t.font.size = Pt(17)  # STRICT >= 14PT
        p_t.font.bold = True
        p_t.font.color.rgb = PRIMARY_TITLE
        p_t.space_after = Pt(3)

        p_sub = tf_t.add_paragraph()
        p_sub.text = data["subtitle"]
        p_sub.font.size = Pt(14)  # STRICT >= 14PT
        p_sub.font.bold = True
        p_sub.font.color.rgb = SUBTITLE_COLOR

        # Image presence check
        has_image = bool(data.get("image_path") and os.path.exists(data["image_path"]))
        if has_image:
            box1_w = Inches(4.9)
            box2_w = Inches(4.7)
            box2_left = Inches(5.55)
            img_left = Inches(10.4)
            img_w = Inches(2.43)
        else:
            box1_w = Inches(5.95)
            box2_w = Inches(5.95)
            box2_left = Inches(6.88)
            img_left = None

        # 4. Box 1 (Left Container - Pure White Card)
        box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.98), box1_w, Inches(3.8))
        box1.fill.solid()
        box1.fill.fore_color.rgb = CARD_BG
        box1.line.color.rgb = CARD_BORDER
        box1.line.width = Pt(1.5)
        tf_1 = box1.text_frame
        tf_1.word_wrap = True
        tf_1.margin_left = Inches(0.18)
        tf_1.margin_top = Inches(0.14)
        tf_1.margin_right = Inches(0.18)
        tf_1.margin_bottom = Inches(0.14)
        
        p1_h = tf_1.paragraphs[0]
        p1_h.text = data["box1_title"]
        p1_h.font.size = Pt(15)  # STRICT >= 14PT
        p1_h.font.bold = True
        p1_h.font.color.rgb = BOX_HEADER_COLOR
        p1_h.space_after = Pt(5)

        if "box1_bullets" in data:
            for bullet in data["box1_bullets"]:
                p_b = tf_1.add_paragraph()
                p_b.text = f"•  {bullet}"
                p_b.font.size = Pt(14)  # STRICT >= 14PT
                p_b.font.color.rgb = BODY_TEXT_COLOR
                p_b.space_after = Pt(4)

        # 5. Box 2 (Middle Container - Pure White Card)
        box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box2_left, Inches(1.98), box2_w, Inches(3.8))
        box2.fill.solid()
        box2.fill.fore_color.rgb = CARD_BG
        box2.line.color.rgb = CARD_BORDER
        box2.line.width = Pt(1.5)
        tf_2 = box2.text_frame
        tf_2.word_wrap = True
        tf_2.margin_left = Inches(0.18)
        tf_2.margin_top = Inches(0.14)
        tf_2.margin_right = Inches(0.18)
        tf_2.margin_bottom = Inches(0.14)

        p2_h = tf_2.paragraphs[0]
        p2_h.text = data["box2_title"]
        p2_h.font.size = Pt(15)  # STRICT >= 14PT
        p2_h.font.bold = True
        p2_h.font.color.rgb = BOX_HEADER_COLOR
        p2_h.space_after = Pt(5)

        if "box2_bullets" in data:
            for bullet in data["box2_bullets"]:
                p_b = tf_2.add_paragraph()
                p_b.text = f"•  {bullet}"
                p_b.font.size = Pt(14)  # STRICT >= 14PT
                p_b.font.color.rgb = BODY_TEXT_COLOR
                p_b.space_after = Pt(4)

        # 5b. Right Image (Mobile App Vector Render)
        if has_image:
            slide.shapes.add_picture(data["image_path"], img_left, Inches(1.98), width=img_w, height=Inches(3.8))

        # 6. Bottom Callout Card (Clean Light Slate Box with Indigo Border)
        bot = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.85), Inches(12.333), Inches(0.85))
        bot.fill.solid()
        bot.fill.fore_color.rgb = CALLOUT_BG
        bot.line.color.rgb = CALLOUT_BORDER
        bot.line.width = Pt(1.5)
        tf_bot = bot.text_frame
        tf_bot.word_wrap = True
        tf_bot.margin_left = Inches(0.18)
        tf_bot.margin_top = Inches(0.08)
        tf_bot.margin_right = Inches(0.18)
        tf_bot.margin_bottom = Inches(0.08)

        pb_h = tf_bot.paragraphs[0]
        pb_h.text = data["bottom_title"]
        pb_h.font.size = Pt(14)  # STRICT >= 14PT
        pb_h.font.bold = True
        pb_h.font.color.rgb = SUBTITLE_COLOR
        pb_h.space_after = Pt(2)

        pb_t = tf_bot.add_paragraph()
        pb_t.text = data["bottom_text"]
        pb_t.font.size = Pt(14)  # STRICT >= 14PT
        pb_t.font.color.rgb = PRIMARY_TITLE
        
        # Add clickable hyperlink for Live App and Figma URL if present in text
        if "app_url" in data:
            run = pb_t.add_run()
            run.text = " [Open Live App]"
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = GREEN_ACCENT
            run.hyperlink.address = data["app_url"]

        if "figma_url" in data:
            run_f = pb_t.add_run()
            run_f.text = " [Open Figma]"
            run_f.font.size = Pt(14)
            run_f.font.bold = True
            run_f.font.color.rgb = SUBTITLE_COLOR
            run_f.hyperlink.address = data["figma_url"]

        # 7. Bottom Navigation Ribbon (Clean Minimalist Pills)
        for idx, step in enumerate(steps_list):
            is_active = (idx + 1) == data["slide_num"]
            step_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5) + Inches(idx * 1.233), Inches(6.82), Inches(1.18), Inches(0.38))
            step_box.fill.solid()
            step_box.fill.fore_color.rgb = SUBTITLE_COLOR if is_active else CALLOUT_BG
            step_box.line.color.rgb = CARD_BORDER
            step_box.line.width = Pt(1)
            tf_s = step_box.text_frame
            tf_s.margin_left = Inches(0)
            tf_s.margin_right = Inches(0)
            tf_s.margin_top = Inches(0.04)
            p_s = tf_s.paragraphs[0]
            p_s.text = step
            p_s.alignment = PP_ALIGN.CENTER
            p_s.font.size = Pt(14)  # STRICT MINIMUM 14PT
            p_s.font.bold = True
            p_s.font.color.rgb = WHITE if is_active else MUTED_TEXT

    out_pptx = "NL Zepto.pptx"
    prs.save(out_pptx)
    print(f"✅ Successfully created ultra-clean PowerPoint presentation: {out_pptx}")
    return slides_data

def build_pdf_deck(slides_data):
    from reportlab.lib.pagesizes import landscape
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

    pdf_filename = "NL Zepto.pdf"
    
    # 16:9 Landscape dimensions (11 x 6.1875 inches = 792 x 445.5 points)
    doc = SimpleDocTemplate(
        pdf_filename,
        pagesize=(792, 445.5),
        rightMargin=20,
        leftMargin=20,
        topMargin=15,
        bottomMargin=15
    )

    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle(
        'SlideTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=14,
        leading=16,
        textColor=colors.HexColor('#0F172A'),
        spaceAfter=2
    )

    subtitle_style = ParagraphStyle(
        'SlideSubtitle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=11,
        leading=13,
        textColor=colors.HexColor('#4F46E5'),
        spaceAfter=4
    )

    tagline_style = ParagraphStyle(
        'TaglineBanner',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=10,
        textColor=colors.HexColor('#FFFFFF')
    )

    h_box_style = ParagraphStyle(
        'BoxHeader',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=11,
        leading=13,
        textColor=colors.HexColor('#4338CA'),
        spaceAfter=4
    )

    body_style = ParagraphStyle(
        'BoxBody',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9.5,
        leading=11.5,
        textColor=colors.HexColor('#1E293B'),
        spaceAfter=3
    )

    bot_title_style = ParagraphStyle(
        'BotTitle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=10,
        leading=12,
        textColor=colors.HexColor('#4F46E5'),
        spaceAfter=1
    )

    bot_text_style = ParagraphStyle(
        'BotText',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9,
        leading=11,
        textColor=colors.HexColor('#0F172A')
    )

    story = []

    for s_idx, data in enumerate(slides_data):
        # 1. Top Banner Tagline Table
        banner_p = Paragraph(f"<b>  {data['tagline'].upper()}</b>", tagline_style)
        slide_num_p = Paragraph(f"<font color='#FFFFFF'><b>SLIDE {data['slide_num']} / 10</b></font>", tagline_style)
        banner_table = Table([[banner_p, slide_num_p]], colWidths=[610, 142])
        banner_table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#3B0764')),
            ('PADDING', (0,0), (-1,-1), 3.5),
            ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
            ('ALIGN', (1,0), (1,0), 'RIGHT')
        ]))
        story.append(banner_table)
        story.append(Spacer(1, 4))

        # 2. Title & Subtitle
        story.append(Paragraph(data['title'], title_style))
        story.append(Paragraph(data['subtitle'], subtitle_style))

        # 3. Two-Column Container Box Grid (Clean Pure White Cards)
        b1_content = [Paragraph(data['box1_title'], h_box_style)]
        if "box1_bullets" in data:
            for bul in data["box1_bullets"]:
                b1_content.append(Paragraph(f"•  {bul}", body_style))

        b2_content = [Paragraph(data['box2_title'], h_box_style)]
        if "box2_bullets" in data:
            for bul in data["box2_bullets"]:
                b2_content.append(Paragraph(f"•  {bul}", body_style))

        if data.get("image_path") and os.path.exists(data["image_path"]):
            img_content = [Image(data["image_path"], width=130, height=210)]
            grid_table = Table([[b1_content, b2_content, img_content]], colWidths=[305, 305, 142])
            grid_table.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#FFFFFF')),
                ('BOX', (0,0), (0,0), 1, colors.HexColor('#CBD5E1')),
                ('BOX', (1,0), (1,0), 1, colors.HexColor('#CBD5E1')),
                ('BOX', (2,0), (2,0), 1, colors.HexColor('#CBD5E1')),
                ('PADDING', (0,0), (-1,-1), 4),
                ('VALIGN', (0,0), (-1,-1), 'TOP'),
                ('ALIGN', (2,0), (2,0), 'CENTER')
            ]))
        else:
            grid_table = Table([[b1_content, b2_content]], colWidths=[371, 371])
            grid_table.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#FFFFFF')),
                ('BOX', (0,0), (0,0), 1, colors.HexColor('#CBD5E1')),
                ('BOX', (1,0), (1,0), 1, colors.HexColor('#CBD5E1')),
                ('PADDING', (0,0), (-1,-1), 5),
                ('VALIGN', (0,0), (-1,-1), 'TOP')
            ]))

        story.append(grid_table)
        story.append(Spacer(1, 4))

        # 4. Bottom Callout Card (Soft Light Tint Card)
        bot_content = [
            Paragraph(data['bottom_title'], bot_title_style),
            Paragraph(data['bottom_text'], bot_text_style)
        ]
        bot_table = Table([[bot_content]], colWidths=[752])
        bot_table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#F1F5F9')),
            ('BOX', (0,0), (-1,-1), 1, colors.HexColor('#4F46E5')),
            ('PADDING', (0,0), (-1,-1), 4),
            ('VALIGN', (0,0), (-1,-1), 'TOP')
        ]))
        story.append(bot_table)

    def draw_bg(canvas, doc):
        canvas.saveState()
        canvas.setFillColor(colors.HexColor('#F8FAFC'))
        canvas.rect(0, 0, 792, 445.5, fill=1, stroke=0)
        canvas.restoreState()

    doc.build(story, onFirstPage=draw_bg, onLaterPages=draw_bg)
    print(f"✅ Successfully created ultra-clean PDF presentation: {pdf_filename}")

if __name__ == "__main__":
    data = build_pptx_deck()
    build_pdf_deck(data)
