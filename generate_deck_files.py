import os
import sys
import shutil

# Force stdout UTF-8 encoding for Windows terminal printing
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def build_pptx_deck():
    prs = Presentation()
    # Widescreen 16:9 (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette (Executive Slate & Zepto Brand Theme)
    BG_DARK = RGBColor(15, 23, 42)          # #0f172a Slate Dark Executive
    CARD_BG = RGBColor(30, 41, 59)          # #1e293b Card container
    CARD_BORDER = RGBColor(99, 102, 241)    # #6366f1 Indigo Accent Border
    HEADER_GRAD = RGBColor(124, 58, 237)    # #7c3aed Zepto Brand Purple Header
    YELLOW_ACCENT = RGBColor(245, 158, 11)   # #f59e0b Amber Gold
    WHITE = RGBColor(248, 250, 252)
    LIGHT_PURPLE = RGBColor(226, 232, 240)
    MUTED_TEXT = RGBColor(148, 163, 184)

    steps_list = ["Context", "Market", "Research", "Insights", "Canvas", "Ideation", "MVP", "Architecture", "Metrics", "GTM"]

    # Image paths (100% Crisp Vector SVG Renders)
    img_dir = os.path.join(os.path.dirname(__file__), 'images', 'crisp_vector_renders')
    img_fallback = os.path.join(os.path.dirname(__file__), 'images', 'tight_iphone_app.png')
    
    def get_img_path(fname):
        p = os.path.join(img_dir, fname)
        return p if os.path.exists(p) else img_fallback

    img1 = get_img_path('01_b2b_free_sampler_cart.png')
    img2 = get_img_path('02_category_streak_quest_board.png')
    img3 = get_img_path('03_restock_cadence_interception.png')
    img4 = get_img_path('04_post_trial_outcome_loop.png')
    img5 = get_img_path('05_lifecycle_moment_interceptor.png')
    img6 = get_img_path('06_freshness_guaranteed_audit.png')
    img7 = get_img_path('07_neighbourhood_trend_feed.png')
    img8 = get_img_path('08_discovery100_voucher_conversion.png')

    slides_data = [
        {
            "slide_num": 1,
            "tagline": "CATEGORY STREAK + GROCERY POINTS CROSS-SUBSIDY MODEL",
            "title": "1. Zepto Expands Blended Gross Margin by Converting Daily Grocery Habits into Recurring Multi-Category Buying",
            "subtitle": "Funding grocery point rewards from high-margin non-grocery categories to bypass traditional P&L subsidy limits.",
            "box1_title": "📋 STRATEGIC BRIEF & PM SCOPE",
            "box1_bullets": [
                "Role & Scope: PM on Zepto Growth Team driving trust-led cross-category discovery.",
                "Routine Tunnel Vision: 71.2% of users are locked into fast daily grocery staple checkouts (<45 seconds).",
                "Strategic Objective: Lift Monthly Active Customers (MAC) buying from 2+ categories from 8.2% to 28.4% in 12 months.",
                "Target Categories: Personal Care (35% margin), Pet Supplies (45% margin), Electronics (30% margin), Baby Care (35% margin)."
            ],
            "box2_title": "⚡ THE CROSS-SUBSIDY ARBITRAGE MODEL",
            "box2_bullets": [
                "Grocery Point Cross-Subsidy: Earning 2x points on morning staples (milk, eggs) offsets the customer's grocery bill.",
                "Sustainable Arbitrage: Points are funded from high non-grocery margins (~35-45%), not from direct P&L discounts.",
                "Ecosystem Lock-in: Rivals cannot copy this cross-subsidy without gross margin arbitrage (or subsidizing from P&L).",
                "Compounding Flywheel: More category exploration -> more grocery points -> stronger grocery habit retention."
            ],
            "image_path": img1,
            "bottom_title": "🔗 PORTAL DIRECTORY & VERIFIED DATA",
            "bottom_text": "Live Streamlit App: https://zeptomvp2.streamlit.app | Public Source Code: NL_Zepto_Growth_PM_Graduation_Project"
        },
        {
            "slide_num": 2,
            "tagline": "CONVENIENCE-JUSTIFIED TOP-UPS VS PLANNED BULK BUYS",
            "title": "2. Users Leak Planned Non-Grocery Spending to Amazon/DMart; MVP Targets Restock Cadence and Trial-to-Repeat Loops",
            "subtitle": "Shifting focus from massive bulk orders to convenience-justified trials and replenishment intercepts.",
            "box1_title": "⭐ WHAT THE ECOSYSTEM HAS SOLVED",
            "box1_bullets": [
                "10-Minute Hyper-Local Delivery Speed (Zepto, Blinkit, Instamart).",
                "Dark Store Density & Real-Time Stock (500+ Dark Stores).",
                "Sub-Second Cart Search & Auto-Complete UI."
            ],
            "box2_title": "❌ CADENCE & TRIAL-TO-REPEAT BLIND SPOTS",
            "box2_bullets": [
                "Purchasing Leakage: 83% of surveyed users (20/24) buy planned non-grocery outside QCommerce on Amazon/Nykaa/DMart [Survey, N=24].",
                "Cadence Suggestion Deficit: Zepto currently has no reorder suggestions based on user time gaps or past items.",
                "Trial-to-Repeat Loop Gap: Users claim brand-sponsored free trial samples once and fail to buy full-size.",
                "Trust & Quality Friction: 20.1% of reviews fear heat degrades dark store items; 17% (4/24) fear refund bot loops."
            ],
            "image_path": img6,
            "bottom_title": "★ CORE THESIS TESTED & VALIDATED",
            "bottom_text": "Category exploration stalls because trial samples are treated as one-offs. To drive repeat, trials must feed a Category Streak Quest that directly offsets daily grocery bills."
        },
        {
            "slide_num": 3,
            "tagline": "AI REVIEW WORKFLOW & PRIMARY RESEARCH SYNTHESIS",
            "title": "3. Multi-Platform Mixed Review Corpus (10k Items) & Survey Isolate 3 Core Pillars: Delights, Feature Requests & Barriers",
            "subtitle": "Synthesizing 10,000 mixed reviews (40% Positive, 30% Neutral, 30% Negative) and primary survey into actionable insights.",
            "box1_title": "📊 AI PIPELINE & MIXED CORPUS (10,000 ITEMS)",
            "box1_bullets": [
                "AI Pipeline: Ingested 10,000 records across 10 channels -> Categorized into Positive Delights (40%), Neutral Feature Requests (30%), Negative Frictions (30%).",
                "40% Positive Delights (4,000 items): High adoption for sub-10 min emergency refills (sunscreen, travel grooming, pet treats, fast phone chargers).",
                "30% Neutral Feature Requests (3,000 items): Unmet demand for AI Shade Matchers, Doorstep Try & Inspect, and Category Streak Points."
            ],
            "box2_title": "🎯 3 SYNTHESIZED BARRIER BUCKETS (USER SURVEY, N=24)",
            "box2_bullets": [
                "1. Mismatch Barrier: Users plan bulk buys elsewhere (33%, 8/24) [Survey Q4]; we intercept them 2 days prior via restock cadence data.",
                "2. Trial-to-Repeat Gap: B2C sampling exists but has no outcome loop or social proof follow-up to convert trial to repeat.",
                "3. Cross-Subsidy Appeal: 63% prefer grocery points on streaks (15/24) over flat cashbacks, locking in daily grocery routines."
            ],
            "image_path": img8,
            "bottom_title": "💡 RESEARCH CONCLUSION & INSIGHT SYNTHESIS",
            "bottom_text": "Users don't need generic ads. They need Category Streaks to lock in loyalty, Cadence Nudges to intercept bulk buys, and Outcome Loop Cards to drive repeat."
        },
        {
            "slide_num": 4,
            "tagline": "HIGH-FREQUENCY GROCERY BUYERS WHO LEAK NON-GROCERY SPEND",
            "title": "4. High-Frequency Grocery Refillers (79% Weekly+) Form the Prime Wedge to Intercept Planned Non-Grocery Restocks",
            "subtitle": "Targeting retained grocery refuelers and using grocery cadence to intercept Amazon purchases 2 days before restock.",
            "box1_title": "📊 TARGET COHORT PROFILE & RATIONALE",
            "box1_bullets": [
                "Cohort Definition: High-frequency grocery users (79% order weekly+, 19/24) who leak non-grocery spend to Amazon/Nykaa (83%, 20/24).",
                "Strategic Fit: Retained users with high basket frequency; lower CAC than acquiring new users.",
                "Predictive Restock Interception: Grocery restock cadence predicts non-grocery replenishment timing, triggering prompts 2 days before Amazon order."
            ],
            "box2_title": "👤 TARGET PERSONA: NEHA (SKINCARE FAN)",
            "box2_bullets": [
                "Profile: Neha, 26, Bangalore · Digital Marketer & Skincare Enthusiast.",
                "JTBD: When ordering my morning milk, I want to try premium skincare samples and lock in cheaper daily groceries, so that category discovery directly offsets my living costs.",
                "Quote: \"If trying a new skincare brand on Zepto gets me 2x points on my daily bread, it becomes a monthly category habit.\"",
                "Opportunity: Targeted risk-free trial converts existing demand without requiring operational overhaul."
            ],
            "image_path": img3,
            "bottom_title": "🎯 IMPACT IF SOLVED FOR THE CAUTIOUS EXPLORER",
            "bottom_text": "Cuts pre-checkout anxiety loops. Makes multi-category exploration safe. Boosts customer LTV by 3.4x via B2B sampling and loyalty points lock-in."
        },
        {
            "slide_num": 5,
            "tagline": "SUPPLY-SIDE MONETIZATION & UNIT ECONOMICS RIGOR",
            "title": "5. Brand-Funded Cross-Subsidy Yields Positive Unit Economics with Immediate Payback on Cross-Category LTV",
            "subtitle": "Grocery rewards are funded by non-grocery margins (~35-45%), not from Zepto P&L.",
            "box1_title": "📈 TAM / SAM / SOM OPPORTUNITY SIZING",
            "box1_bullets": [
                "TAM: $18.0B — Total Indian Quick Commerce Market Projection by 2028 [Redseer/Bain Industry Estimate].",
                "SAM: $4.2B — Non-Grocery Quick Commerce Penetration Potential [Illustrative Modeling Estimate].",
                "SOM: $480M — Zepto Cross-Category Discovery Capture Opportunity [Illustrative Modeling Estimate]."
            ],
            "box2_title": "💰 UNIT ECONOMICS & CONVERSION LEVERS (N=24)",
            "box2_bullets": [
                "Grocery Streak Hook: Lift stagnant 8.2% MAC to 28.4% MAC exploration via points.",
                "Cross-Subsidy Arbitrage: Funded from non-grocery margins (35-45%), keeping points loop sustainable.",
                "B2B Brand Samples: Brand funded at Rs. 15/unit, acting as a zero-CAC trial on-ramp in bag."
            ],
            "image_path": None,
            "bottom_title": "💡 FINANCIAL FLYWHEEL SUMMARY",
            "bottom_text": "Shifting grocery refillers to 35%–50% margin categories expands gross margin by +300bps and captures $480M SOM."
        },
        {
            "slide_num": 6,
            "tagline": "COMPOUNDING DATA FLYWHEEL VS COPIABLE PROMOS",
            "title": "6. Zepto's Defensibility Lies in its Restock Cadence Engine and Local Purchase-Graph, Not Copiable Promos",
            "subtitle": "Flat discounts are easily copied; restock history and PIN-code purchase graphs create years of defensibility.",
            "box1_title": "⚡ TRIVIALLY COPIABLE VS COMPOUNDING MOATS",
            "box1_bullets": [
                "Trivially Copiable Promos: Free samples, coupons, and 15-min swap policies can be copied by Blinkit/Instamart in a week.",
                "Compounding Data Moat 1: Requires 6-12 months per-user order history to generate accurate restock timing.",
                "Compounding Brand Moat 2: Hyperlocal social proof requires years of PIN-code level dark-store purchase graphs."
            ],
            "box2_title": "🏆 THREE STRATEGIC HORIZONS (RICE MATRIX)",
            "box2_bullets": [
                "Horizon 1 (MVP) — Streak Loop, Cadence Intercept & Outcome Cards: Habitual lock-in & trial-to-repeat conversion. RICE: 230.0 [Vetted MVP].",
                "Horizon 2 (Growth) — Lifecycle Interceptor & Neighbourhood Trends: Basket-trigger discovery & local proof. RICE: 190.0 [Horizon 2].",
                "Horizon 3 (Vision) — Predictive Subscription Auto-Replenishment: 100% automated smart grocery companion [RICE: 150.0]."
            ],
            "image_path": None,
            "bottom_title": "💡 WHY HORIZON 1 (MVP) WINS FIRST",
            "bottom_text": "B2B Sampling requires zero capital expenditure (brands fund sample units) and addresses the core barrier (risk) right inside recurring grocery bags."
        },
        {
            "slide_num": 7,
            "tagline": "THE FLYWHEEL MVP: STREAKS, CADENCE NUDGES, OUTCOME CARDS",
            "title": "7. The Discovery MVP Embeds Category Streaks, Cadence Nudges, and Post-Trial Outcome Loop Cards",
            "subtitle": "Connecting trials, streaks, and intercepts into a closed-loop category adoption system.",
            "box1_title": "💎 BUILT MVP CAPABILITIES 1, 2 & 3",
            "box1_bullets": [
                "1. 🏆 Category Streak Board (Endowed Progress): Try 1 new category this month -> earn 2x grocery points. Features pre-stamped head start.",
                "2. 🕒 Cadence Interception Nudge: Fires nudge 2 days before Amazon bulk buy (offers ₹0 brand-funded trial pack).",
                "3. 📦 Outcome Loop Card: Post-trial repeat card. Fires at 48 hours with PIN-code proof ('847 near you reordered Cetaphil'). Rate & Add to next cart."
            ],
            "box2_title": "🏆 HORIZON 2 EXPANSIONS & TRUST BADGES",
            "box2_bullets": [
                "4. 🔄 Lifecycle Moment Interceptor: Cart diaper addition prompts user declaration, unlocking curated discovery trail.",
                "5. 🔥 Neighbourhood Trend Feed: Surfaces aggregated Sector 56 trends with strict 50+ user privacy filters.",
                "6. 🛡️ Freshness Guaranteed Badge: Replaces costly CCTV sensors with simple checkout trust badge and doorstep audit logs."
            ],
            "image_path": img2,
            "bottom_title": "🎨 LIVE FIGMA DESIGN SYSTEM & INTERACTIVE MVP PORTAL",
            "bottom_text": "Figma Design & Wireframes: https://www.figma.com/design/2gZJHtjmpnI677IhC66PqZ/Untitled?node-id=0-1 | Live App: https://zeptomvp2.streamlit.app"
        },
        {
            "slide_num": 8,
            "tagline": "FROM A TYPED GROCERY CART TO A TRUSTED CROSS-CATEGORY ORDER. SAME APP, TWO VIEWS.",
            "title": "8. Four-Layer Decision Engine Powers Predictive Restock Nudges and Closed-Loop Brand Attribution",
            "subtitle": "Four-layer architecture integrates brand inventory, cart triggers, and loyalty ledgers.",
            "box1_title": "⚙ SYSTEM ARCHITECTURE (4 CORE LAYERS)",
            "box1_bullets": [
                "Layer 1 (Client UI): Mobile Cart UI, pre-stamped Category Streak Board, Cadence restock widget, Outcome Loop Card.",
                "Layer 2 (Decision Engine): Restock Cadence Inference Engine, Endowed Progress Tracker, Hyperlocal Social Graph Parser.",
                "Layer 3 (Operations): Dark-store picker checklist update for B2B brand sample insertions (under 3s packing SLA).",
                "Layer 4 (B2B Marketplace Portal): Brand attribution dashboard tracking sample-to-repeat conversion rates for FMCG brands."
            ],
            "box2_title": "🔄 EMOTION & METRIC MAPPING ACROSS STAGES",
            "box2_bullets": [
                "Stage 1 (Cart): Types grocery staples -> System flags pre-stamped streak and restock cadence -> Curious",
                "Stage 2 (Nudge): Claims free trial sample -> System packs sample directly in grocery bag -> Confident",
                "Stage 3 (Audit): Views freshness guarantee log -> System verifies climate logs -> Reassured",
                "Stage 4 (Checkout): Receives outcome loop follow-up -> Rates & adds full-size to cart -> Empowered"
            ],
            "image_path": img4,
            "bottom_title": "💡 TECHNICAL ARCHITECTURE MOAT & FEASIBILITY",
            "bottom_text": "Zero hardware cost for MVP sampler. Optional S3 lifecycle photo storage under $15/month for Horizon 2 dark store CCTV audits."
        },
        {
            "slide_num": 9,
            "tagline": "INCREMENTAL LIFT VIA RANDOMIZED HOLDOUT GROUPS",
            "title": "9. MCER Growth (8.2% -> 28.4%) Is Validated via Randomized Holdout Groups and SLA Guardrails",
            "subtitle": "Measuring true incremental lift via control groups rather than raw attach proxies.",
            "box1_title": "⭐ NORTH STAR & INCREMENTALITY FRAMEWORK",
            "box1_bullets": [
                "North Star Metric: Monthly Category Exploration Rate (MCER) — % of MACs purchasing from 2+ categories/month.",
                "• Baseline: 8.2% MAC -> Target: 28.4% MAC in 12 months [Illustrative Trajectory Target].",
                "Randomized Holdout Control: 10% control group (no sample/nudge) vs treatment to measure true incremental causal lift.",
                "Leading Indicators: Sample attach rate (>35%), Sample-to-repeat conversion (12% in 14 days), Category streak completion rate."
            ],
            "box2_title": "🛡️ OPERATIONAL GUARDRAIL METRICS",
            "box2_bullets": [
                "Picker SLA Floor: Dark-store sample packing time addition capped at <3 seconds per order.",
                "Checkout Drop-off Cap: Nudge interaction must not increase cart drop-off rate (>0.2%).",
                "Refund Rate Ceiling: Doorstep replacement requests must remain below 1.5% of non-grocery orders.",
                "Privacy Guardrails: Neighbourhood feed requires 50+ users PIN-code threshold and sensitive category blocks."
            ],
            "image_path": img7,
            "bottom_title": "📈 METRIC COMPOUNDING & INTEGRITY",
            "bottom_text": "MCER is measured strictly from real event streams (sample claims, voucher redemptions, category streak completions). Zero proxies."
        },
        {
            "slide_num": 10,
            "tagline": "HORIZON 2 ROADMAP & PRIVACY GUARDRAILS",
            "title": "10. Phased GTM Roadmap Launches Horizon 1 MVP While Building Horizon 2 Privacy-Guardrailed Features",
            "subtitle": "Phased rollout for Category Streaks, Cadence Intercepts, Lifecycle Moments, and Trend Feeds.",
            "box1_title": "🚀 3-PHASE ROLLOUT & PRIVACY SEGREGATION",
            "box1_bullets": [
                "Phase 1 (MVP Beta): 30-day pilot of Streaks + Cadence Intercepts + Outcome Cards in Bangalore (10 dark stores).",
                "Phase 2 (Growth Rollout): Expansion of Horizon 1 to all Metro cities (150 dark stores); launch brand B2B portal.",
                "Phase 3 (Horizon 2 Beta): Rollout of Lifecycle Moment Interceptor & Neighbourhood Trend Feed under strict privacy rules.",
                "Phase 4 (GA Rollout): Full national rollout across 500+ dark store hubs."
            ],
            "box2_title": "⚠️ RISKS & MITIGATIONS SUMMARY",
            "box2_bullets": [
                "Rival Copying: Ecosystem points streak lock-in on daily staples (requires margin structure).",
                "Data Runway: Requires 6-12 months per-user restock cadence order history to tune model.",
                "Privacy Backlash: Opt-in only neighborhood feeds, PIN-code minimum, and 50+ user count thresholds."
            ],
            "image_path": img5,
            "bottom_title": "🔗 ANONYMOUS REVIEWER DATA & DEMO DIRECTORY",
            "bottom_text": "Figma Design & Wireframes: https://www.figma.com/design/2gZJHtjmpnI677IhC66PqZ/Untitled?node-id=0-1 | Live Streamlit: https://zeptomvp2.streamlit.app"
        }
    ]

    for data in slides_data:
        slide = prs.slides.add_slide(blank_layout)

        # 1. Background Fill
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = BG_DARK
        bg_shape.line.color.rgb = BG_DARK

        # 2. Top Header Tagline Banner
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.35), Inches(12.333), Inches(0.5))
        banner.fill.solid()
        banner.fill.fore_color.rgb = HEADER_GRAD
        banner.line.color.rgb = HEADER_GRAD
        tf_b = banner.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = Inches(0.15)
        tf_b.margin_top = Inches(0.08)
        p_b = tf_b.paragraphs[0]
        p_b.text = f"  {data['tagline'].upper()}"
        p_b.font.size = Pt(14)
        p_b.font.bold = True
        p_b.font.color.rgb = WHITE

        # 3. Slide Title & Subtitle
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.90), Inches(12.333), Inches(1.05))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = Inches(0)
        tf_t.margin_top = Inches(0)
        
        p_t = tf_t.paragraphs[0]
        p_t.text = data["title"]
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = WHITE
        p_t.space_after = Pt(2)

        p_sub = tf_t.add_paragraph()
        p_sub.text = data["subtitle"]
        p_sub.font.size = Pt(14)
        p_sub.font.color.rgb = YELLOW_ACCENT

        # Image presence check (Slide 4 and Slide 7)
        has_image = bool(data.get("image_path") and os.path.exists(data["image_path"]))
        if has_image:
            box1_w = Inches(4.9)
            box2_w = Inches(4.7)
            box2_left = Inches(5.55)
            img_left = Inches(10.4)
            img_w = Inches(2.43)
        else:
            box1_w = Inches(5.95)
            box2_w = Inches(5.95)
            box2_left = Inches(6.88)
            img_left = None

        # 4. Box 1 (Left Container)
        box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.98), box1_w, Inches(3.8))
        box1.fill.solid()
        box1.fill.fore_color.rgb = CARD_BG
        box1.line.color.rgb = CARD_BORDER
        box1.line.width = Pt(1.5)
        tf_1 = box1.text_frame
        tf_1.word_wrap = True
        tf_1.margin_left = Inches(0.15)
        tf_1.margin_top = Inches(0.12)
        tf_1.margin_right = Inches(0.15)
        tf_1.margin_bottom = Inches(0.12)
        
        p1_h = tf_1.paragraphs[0]
        p1_h.text = data["box1_title"]
        p1_h.font.size = Pt(15)
        p1_h.font.bold = True
        p1_h.font.color.rgb = YELLOW_ACCENT
        p1_h.space_after = Pt(4)

        if "box1_text" in data:
            p1_b = tf_1.add_paragraph()
            p1_b.text = data["box1_text"]
            p1_b.font.size = Pt(14)
            p1_b.font.color.rgb = LIGHT_PURPLE
            p1_b.space_after = Pt(4)
        elif "box1_bullets" in data:
            for bullet in data["box1_bullets"]:
                p_b = tf_1.add_paragraph()
                p_b.text = f"•  {bullet}"
                p_b.font.size = Pt(14)
                p_b.font.color.rgb = LIGHT_PURPLE
                p_b.space_after = Pt(4)

        # 5. Box 2 (Middle Container)
        box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box2_left, Inches(1.98), box2_w, Inches(3.8))
        box2.fill.solid()
        box2.fill.fore_color.rgb = CARD_BG
        box2.line.color.rgb = CARD_BORDER
        box2.line.width = Pt(1.5)
        tf_2 = box2.text_frame
        tf_2.word_wrap = True
        tf_2.margin_left = Inches(0.15)
        tf_2.margin_top = Inches(0.12)
        tf_2.margin_right = Inches(0.15)
        tf_2.margin_bottom = Inches(0.12)

        p2_h = tf_2.paragraphs[0]
        p2_h.text = data["box2_title"]
        p2_h.font.size = Pt(15)
        p2_h.font.bold = True
        p2_h.font.color.rgb = YELLOW_ACCENT
        p2_h.space_after = Pt(4)

        if "box2_text" in data:
            p2_b = tf_2.add_paragraph()
            p2_b.text = data["box2_text"]
            p2_b.font.size = Pt(14)
            p2_b.font.color.rgb = LIGHT_PURPLE
            p2_b.space_after = Pt(4)
        elif "box2_bullets" in data:
            for bullet in data["box2_bullets"]:
                p_b = tf_2.add_paragraph()
                p_b.text = f"•  {bullet}"
                p_b.font.size = Pt(14)
                p_b.font.color.rgb = LIGHT_PURPLE
                p_b.space_after = Pt(4)

        # 5b. Right Image (Mobile Phone App Mockup)
        if has_image:
            slide.shapes.add_picture(data["image_path"], img_left, Inches(1.98), width=img_w, height=Inches(3.8))

        # 6. Bottom Callout Card
        bot = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.85), Inches(12.333), Inches(0.85))
        bot.fill.solid()
        bot.fill.fore_color.rgb = CARD_BG
        bot.line.color.rgb = YELLOW_ACCENT
        bot.line.width = Pt(1.5)
        tf_bot = bot.text_frame
        tf_bot.word_wrap = True
        tf_bot.margin_left = Inches(0.18)
        tf_bot.margin_top = Inches(0.08)
        tf_bot.margin_right = Inches(0.18)
        tf_bot.margin_bottom = Inches(0.08)

        pb_h = tf_bot.paragraphs[0]
        pb_h.text = data["bottom_title"]
        pb_h.font.size = Pt(14)
        pb_h.font.bold = True
        pb_h.font.color.rgb = YELLOW_ACCENT
        pb_h.space_after = Pt(2)

        pb_t = tf_bot.add_paragraph()
        pb_t.text = data["bottom_text"]
        pb_t.font.size = Pt(14)
        pb_t.font.color.rgb = WHITE

        # 7. Bottom Navigation Ribbon (Strict 14pt minimum font size)
        for idx, step in enumerate(steps_list):
            is_active = (idx + 1) == data["slide_num"]
            step_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5) + Inches(idx * 1.233), Inches(6.82), Inches(1.18), Inches(0.38))
            step_box.fill.solid()
            step_box.fill.fore_color.rgb = YELLOW_ACCENT if is_active else CARD_BG
            step_box.line.color.rgb = CARD_BORDER
            step_box.line.width = Pt(1)
            tf_s = step_box.text_frame
            tf_s.margin_left = Inches(0)
            tf_s.margin_right = Inches(0)
            tf_s.margin_top = Inches(0.04)
            p_s = tf_s.paragraphs[0]
            p_s.text = step
            p_s.alignment = PP_ALIGN.CENTER
            p_s.font.size = Pt(14)
            p_s.font.bold = True
            p_s.font.color.rgb = BG_DARK if is_active else MUTED_TEXT

    # Save Fellowship Compliant Filename
    for fname in ["NL_Zepto_Growth_PM_Graduation_Project.pptx", "NL_Zepto_Growth_PM_Graduation_Project_v2.pptx", "NL_Zepto_Growth_PM_Graduation_Project_v3.pptx"]:
        try:
            prs.save(fname)
            out_pptx = fname
            break
        except Exception:
            continue
        
    print(f"✅ Successfully created PowerPoint presentation: {out_pptx}")
    return slides_data

def build_pdf_deck(slides_data):
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, KeepTogether, Image
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

    pdf_filename = "NL_Zepto_Growth_PM_Graduation_Project.pdf"
    
    # 16:9 Landscape dimensions (11 x 6.1875 inches = 792 x 445.5 points)
    doc = SimpleDocTemplate(
        pdf_filename,
        pagesize=(792, 445.5),
        rightMargin=20,
        leftMargin=20,
        topMargin=15,
        bottomMargin=15
    )

    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle(
        'SlideTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=13.5,
        leading=15,
        textColor=colors.HexColor('#FFFFFF'),
        spaceAfter=2
    )

    subtitle_style = ParagraphStyle(
        'SlideSubtitle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=8.5,
        leading=10.5,
        textColor=colors.HexColor('#FFD700'),
        spaceAfter=5
    )

    tagline_style = ParagraphStyle(
        'TaglineBanner',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=9,
        textColor=colors.HexColor('#FFFFFF')
    )

    h_box_style = ParagraphStyle(
        'BoxHeader',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=9.5,
        leading=11.5,
        textColor=colors.HexColor('#FFD700'),
        spaceAfter=4
    )

    body_style = ParagraphStyle(
        'BoxBody',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=7.5,
        leading=9.5,
        textColor=colors.HexColor('#E2D9F3'),
        spaceAfter=3
    )

    bot_title_style = ParagraphStyle(
        'BotTitle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=8.5,
        leading=10.5,
        textColor=colors.HexColor('#FFD700'),
        spaceAfter=1
    )

    bot_text_style = ParagraphStyle(
        'BotText',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=7,
        leading=8.8,
        textColor=colors.HexColor('#FFFFFF')
    )

    steps_list = ["Context", "Market", "Research", "Insights", "Canvas", "Ideation", "MVP", "Architecture", "Metrics", "GTM"]

    story = []

    for s_idx, data in enumerate(slides_data):
        # 1. Top Banner Tagline Table
        banner_p = Paragraph(f"<b>  {data['tagline'].upper()}</b>", tagline_style)
        slide_num_p = Paragraph(f"<font color='#FFFFFF'><b>SLIDE {data['slide_num']} / 10</b></font>", tagline_style)
        banner_table = Table([[banner_p, slide_num_p]], colWidths=[610, 142])
        banner_table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#E05238')),
            ('PADDING', (0,0), (-1,-1), 3.5),
            ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
            ('ALIGN', (1,0), (1,0), 'RIGHT')
        ]))
        story.append(banner_table)
        story.append(Spacer(1, 4))

        # 2. Title & Subtitle
        story.append(Paragraph(data['title'], title_style))
        story.append(Paragraph(data['subtitle'], subtitle_style))

        # 3. Two-Column Container Box Grid
        # Box 1 content
        b1_content = [Paragraph(data['box1_title'], h_box_style)]
        if "box1_text" in data:
            b1_content.append(Paragraph(data['box1_text'], body_style))
        elif "box1_bullets" in data:
            for bul in data["box1_bullets"]:
                b1_content.append(Paragraph(f"•  {bul}", body_style))

        # Box 2 content
        b2_content = [Paragraph(data['box2_title'], h_box_style)]
        if "box2_text" in data:
            b2_content.append(Paragraph(data['box2_text'], body_style))
        elif "box2_bullets" in data:
            for bul in data["box2_bullets"]:
                b2_content.append(Paragraph(f"•  {bul}", body_style))

        if data.get("image_path") and os.path.exists(data["image_path"]):
            img_content = [Image(data["image_path"], width=130, height=210)]
            grid_table = Table([[b1_content, b2_content, img_content]], colWidths=[305, 305, 142])
            grid_table.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#1E1430')),
                ('BOX', (0,0), (0,0), 1, colors.HexColor('#8224E3')),
                ('BOX', (1,0), (1,0), 1, colors.HexColor('#8224E3')),
                ('BOX', (2,0), (2,0), 1, colors.HexColor('#8224E3')),
                ('PADDING', (0,0), (-1,-1), 4),
                ('VALIGN', (0,0), (-1,-1), 'TOP'),
                ('ALIGN', (2,0), (2,0), 'CENTER')
            ]))
        else:
            grid_table = Table([[b1_content, b2_content]], colWidths=[371, 371])
            grid_table.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#1E1430')),
                ('BOX', (0,0), (0,0), 1, colors.HexColor('#8224E3')),
                ('BOX', (1,0), (1,0), 1, colors.HexColor('#8224E3')),
                ('PADDING', (0,0), (-1,-1), 5),
                ('VALIGN', (0,0), (-1,-1), 'TOP')
            ]))

        story.append(grid_table)
        story.append(Spacer(1, 4))

        # 4. Bottom Callout Card
        bot_content = [
            Paragraph(data['bottom_title'], bot_title_style),
            Paragraph(data['bottom_text'], bot_text_style)
        ]
        bot_table = Table([[bot_content]], colWidths=[752])
        bot_table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#1E1430')),
            ('BOX', (0,0), (-1,-1), 1, colors.HexColor('#FFD700')),
            ('PADDING', (0,0), (-1,-1), 4)
        ]))
        story.append(bot_table)
        story.append(Spacer(1, 4))

        # 5. Bottom 10-Step Nav Bar
        nav_cells = []
        for n_idx, st_name in enumerate(steps_list):
            is_act = (n_idx + 1) == data['slide_num']
            c_fg = '#000000' if is_act else '#A096B4'
            p_st = Paragraph(f"<font color='{c_fg}'><b>{st_name}</b></font>", ParagraphStyle('NavSt', fontName='Helvetica-Bold', fontSize=6.5, alignment=1))
            nav_cells.append(p_st)

        nav_table = Table([nav_cells], colWidths=[75.2]*10)
        nav_style = [('VALIGN', (0,0), (-1,-1), 'MIDDLE'), ('PADDING', (0,0), (-1,-1), 2)]
        for n_idx in range(10):
            is_act = (n_idx + 1) == data['slide_num']
            bg_col = colors.HexColor('#FFD700') if is_act else colors.HexColor('#1E1430')
            nav_style.append(('BACKGROUND', (n_idx, 0), (n_idx, 0), bg_col))
            nav_style.append(('BOX', (n_idx, 0), (n_idx, 0), 0.5, colors.HexColor('#8224E3')))
        
        nav_table.setStyle(TableStyle(nav_style))
        story.append(nav_table)

        if s_idx < len(slides_data) - 1:
            story.append(PageBreak())

    def draw_bg(canvas, doc):
        canvas.saveState()
        canvas.setFillColor(colors.HexColor('#130D1E'))
        canvas.rect(0, 0, 792, 445.5, fill=1, stroke=0)
        canvas.restoreState()

    doc.build(story, onFirstPage=draw_bg, onLaterPages=draw_bg)
    print(f"✅ Successfully created PDF presentation: {pdf_filename}")

if __name__ == "__main__":
    try:
        data = build_pptx_deck()
        build_pdf_deck(data)
    except Exception as e:
        print("Error generating deck files:", e)
